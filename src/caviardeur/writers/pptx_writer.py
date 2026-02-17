from pathlib import Path

from pptx import Presentation

from ..readers.base import DocumentContent


def write_pptx(content: DocumentContent, output_path: Path, source_path: Path) -> None:
    """Write pseudonymized content back to a .pptx file.

    Opens the original presentation and replaces run texts in place to preserve formatting.
    """
    prs = Presentation(str(source_path))

    # Build lookup: (location tuple) -> new text
    chunk_map: dict[tuple, str] = {}
    for chunk in content.chunks:
        loc = chunk.location
        loc_type = loc.get("type", "")

        if loc_type == "pptx_run":
            key = ("pptx_run", loc["slide_idx"], loc["shape_idx"], loc["para_idx"], loc["run_idx"])
            chunk_map[key] = chunk.text
        elif loc_type == "pptx_table_run":
            key = (
                "pptx_table_run",
                loc["slide_idx"],
                loc["shape_idx"],
                loc["row_idx"],
                loc["cell_idx"],
                loc["para_idx"],
                loc["run_idx"],
            )
            chunk_map[key] = chunk.text

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(para.runs):
                        key = ("pptx_run", slide_idx, shape_idx, para_idx, run_idx)
                        if key in chunk_map:
                            run.text = chunk_map[key]

            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        for para_idx, para in enumerate(cell.text_frame.paragraphs):
                            for run_idx, run in enumerate(para.runs):
                                key = (
                                    "pptx_table_run",
                                    slide_idx,
                                    shape_idx,
                                    row_idx,
                                    cell_idx,
                                    para_idx,
                                    run_idx,
                                )
                                if key in chunk_map:
                                    run.text = chunk_map[key]

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
