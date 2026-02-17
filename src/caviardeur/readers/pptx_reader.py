from pathlib import Path

from pptx import Presentation

from .base import DocumentContent, TextChunk


def read_pptx(path: Path) -> DocumentContent:
    """Read a .pptx file, extracting text at the run level from shapes and tables."""
    prs = Presentation(str(path))
    chunks: list[TextChunk] = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(para.runs):
                        if run.text:
                            chunks.append(
                                TextChunk(
                                    text=run.text,
                                    location={
                                        "type": "pptx_run",
                                        "slide_idx": slide_idx,
                                        "shape_idx": shape_idx,
                                        "para_idx": para_idx,
                                        "run_idx": run_idx,
                                    },
                                )
                            )
                    chunks.append(
                        TextChunk(
                            text="\n",
                            location={
                                "type": "pptx_separator",
                                "slide_idx": slide_idx,
                                "shape_idx": shape_idx,
                                "para_idx": para_idx,
                            },
                        )
                    )

            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        for para_idx, para in enumerate(cell.text_frame.paragraphs):
                            for run_idx, run in enumerate(para.runs):
                                if run.text:
                                    chunks.append(
                                        TextChunk(
                                            text=run.text,
                                            location={
                                                "type": "pptx_table_run",
                                                "slide_idx": slide_idx,
                                                "shape_idx": shape_idx,
                                                "row_idx": row_idx,
                                                "cell_idx": cell_idx,
                                                "para_idx": para_idx,
                                                "run_idx": run_idx,
                                            },
                                        )
                                    )

        # Separator between slides
        chunks.append(
            TextChunk(
                text="\n",
                location={"type": "pptx_slide_separator", "slide_idx": slide_idx},
            )
        )

    content = DocumentContent(chunks=chunks, metadata={"source_path": str(path), "format": "pptx"})
    content.assign_offsets()
    return content
