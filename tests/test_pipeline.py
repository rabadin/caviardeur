"""Integration tests for the full pipeline (using mocked NER to avoid model download)."""

from pathlib import Path
from unittest.mock import patch

from caviardeur.config import Config
from caviardeur.detectors.base import DetectedEntity, EntityType
from caviardeur.pipeline import process_file
from caviardeur.pseudonymizer.mapping import MappingStore

FIXTURES = Path(__file__).parent / "fixtures"


def _mock_detect_all(text, **kwargs):
    """Mock detector that finds known PII in test text."""
    entities = []
    targets = [
        ("Jean Dupont", EntityType.PERSON, 0.96),
        ("Marie Laurent", EntityType.PERSON, 0.94),
        ("Sophie Bernard", EntityType.PERSON, 0.93),
        ("Nextech Solutions SAS", EntityType.COMPANY, 0.97),
        ("DataFlow Industries", EntityType.COMPANY, 0.95),
        ("Société Exemple SAS", EntityType.COMPANY, 0.87),
        ("73282932000074", EntityType.SIRET, 0.95),
        ("732 829 320 00074", EntityType.SIRET, 0.95),
    ]
    for target_text, etype, conf in targets:
        start = 0
        while True:
            idx = text.find(target_text, start)
            if idx == -1:
                break
            entities.append(
                DetectedEntity(
                    entity_type=etype,
                    text=target_text,
                    start=idx,
                    end=idx + len(target_text),
                    confidence=conf,
                    source="mock",
                )
            )
            start = idx + 1
    return entities


# --- TXT ---


@patch("caviardeur.pipeline.detect_all", side_effect=_mock_detect_all)
def test_pipeline_txt(mock_detect, tmp_path: Path):
    input_dir = tmp_path / "input"
    input_dir.mkdir()
    txt_file = input_dir / "test.txt"
    txt_file.write_text(
        "Jean Dupont travaille chez Société Exemple SAS.\nSIRET: 73282932000074\nContact: Marie Laurent",
        encoding="utf-8",
    )

    output_dir = tmp_path / "output"
    config = Config(output_dir=output_dir, dry_run=False)
    mapping = MappingStore()

    entities = process_file(txt_file, config, mapping)

    assert len(entities) == 4

    output_file = output_dir / "test.txt"
    assert output_file.exists()
    content = output_file.read_text(encoding="utf-8")
    assert "Jean Dupont" not in content
    assert "Marie Laurent" not in content
    assert "Société Exemple SAS" not in content
    assert "73282932000074" not in content
    assert "PERSON_001" in content
    assert "PERSON_002" in content
    assert "COMPANY_001" in content
    assert "SIRET_001" in content


@patch("caviardeur.pipeline.detect_all", side_effect=_mock_detect_all)
def test_pipeline_dry_run(mock_detect, tmp_path: Path):
    input_dir = tmp_path / "input"
    input_dir.mkdir()
    txt_file = input_dir / "test.txt"
    txt_file.write_text("Jean Dupont est ici.", encoding="utf-8")

    output_dir = tmp_path / "output"
    config = Config(output_dir=output_dir, dry_run=True)
    mapping = MappingStore()

    entities = process_file(txt_file, config, mapping)

    assert len(entities) == 1
    assert not output_dir.exists()


@patch("caviardeur.pipeline.detect_all", side_effect=_mock_detect_all)
def test_pipeline_mapping_persistence(mock_detect, tmp_path: Path):
    input_dir = tmp_path / "input"
    input_dir.mkdir()

    f1 = input_dir / "batch1.txt"
    f1.write_text("Jean Dupont est directeur.", encoding="utf-8")

    output1 = tmp_path / "out1"
    config1 = Config(output_dir=output1, dry_run=False)
    mapping = MappingStore()
    process_file(f1, config1, mapping)
    mapping_path = output1 / "mapping.json"
    mapping.save(mapping_path)

    f2 = input_dir / "batch2.txt"
    f2.write_text("Contact: Jean Dupont", encoding="utf-8")

    output2 = tmp_path / "out2"
    config2 = Config(output_dir=output2, dry_run=False)
    mapping2 = MappingStore.load(mapping_path)
    process_file(f2, config2, mapping2)

    content = (output2 / "batch2.txt").read_text(encoding="utf-8")
    assert "PERSON_001" in content


# --- Markdown ---


@patch("caviardeur.pipeline.detect_all", side_effect=_mock_detect_all)
def test_pipeline_md(mock_detect, tmp_path: Path):
    output_dir = tmp_path / "output"
    config = Config(output_dir=output_dir, dry_run=False)
    mapping = MappingStore()

    entities = process_file(FIXTURES / "sample.md", config, mapping)

    assert len(entities) > 0
    output_file = output_dir / "sample.md"
    assert output_file.exists()
    content = output_file.read_text(encoding="utf-8")
    assert "Jean Dupont" not in content
    assert "PERSON_001" in content
    # Markdown structure preserved
    assert "# Compte-rendu" in content


# --- JSON ---


@patch("caviardeur.pipeline.detect_all", side_effect=_mock_detect_all)
def test_pipeline_json(mock_detect, tmp_path: Path):
    output_dir = tmp_path / "output"
    config = Config(output_dir=output_dir, dry_run=False)
    mapping = MappingStore()

    entities = process_file(FIXTURES / "sample.json", config, mapping)

    assert len(entities) > 0
    output_file = output_dir / "sample.json"
    assert output_file.exists()
    content = output_file.read_text(encoding="utf-8")
    assert "Jean Dupont" not in content
    assert "PERSON_001" in content
    # JSON structure preserved (braces still there)
    assert '"nom"' in content


# --- XML ---


@patch("caviardeur.pipeline.detect_all", side_effect=_mock_detect_all)
def test_pipeline_xml(mock_detect, tmp_path: Path):
    output_dir = tmp_path / "output"
    config = Config(output_dir=output_dir, dry_run=False)
    mapping = MappingStore()

    entities = process_file(FIXTURES / "sample.xml", config, mapping)

    assert len(entities) > 0
    output_file = output_dir / "sample.xml"
    assert output_file.exists()
    content = output_file.read_text(encoding="utf-8")
    assert "Jean Dupont" not in content
    assert "PERSON_001" in content
    # XML tags preserved
    assert "<nom>" in content


# --- DOCX (from fixture) ---


@patch("caviardeur.pipeline.detect_all", side_effect=_mock_detect_all)
def test_pipeline_docx(mock_detect, tmp_path: Path):
    output_dir = tmp_path / "output"
    config = Config(output_dir=output_dir, dry_run=False)
    mapping = MappingStore()

    entities = process_file(FIXTURES / "sample.docx", config, mapping)

    assert len(entities) > 0
    output_file = output_dir / "sample.docx"
    assert output_file.exists()


# --- XLSX (from fixture) ---


@patch("caviardeur.pipeline.detect_all", side_effect=_mock_detect_all)
def test_pipeline_xlsx(mock_detect, tmp_path: Path):
    output_dir = tmp_path / "output"
    config = Config(output_dir=output_dir, dry_run=False)
    mapping = MappingStore()

    entities = process_file(FIXTURES / "sample.xlsx", config, mapping)

    assert len(entities) > 0
    output_file = output_dir / "sample.xlsx"
    assert output_file.exists()


# --- PDF (from fixture) ---


@patch("caviardeur.pipeline.detect_all", side_effect=_mock_detect_all)
def test_pipeline_pdf(mock_detect, tmp_path: Path):
    output_dir = tmp_path / "output"
    config = Config(output_dir=output_dir, dry_run=False)
    mapping = MappingStore()

    entities = process_file(FIXTURES / "sample.pdf", config, mapping)

    assert len(entities) > 0
    output_file = output_dir / "sample.pdf"
    assert output_file.exists()

    # Verify the pdf can be opened and PII is redacted
    import fitz

    doc = fitz.open(str(output_file))
    all_text = ""
    for page in doc:
        all_text += page.get_text()
    doc.close()
    assert "Jean Dupont" not in all_text
    assert "PERSON" in all_text


# --- PPTX (from fixture) ---


@patch("caviardeur.pipeline.detect_all", side_effect=_mock_detect_all)
def test_pipeline_pptx(mock_detect, tmp_path: Path):
    output_dir = tmp_path / "output"
    config = Config(output_dir=output_dir, dry_run=False)
    mapping = MappingStore()

    entities = process_file(FIXTURES / "sample.pptx", config, mapping)

    assert len(entities) > 0
    output_file = output_dir / "sample.pptx"
    assert output_file.exists()

    # Verify the pptx can be opened and contains pseudonyms
    from pptx import Presentation

    prs = Presentation(str(output_file))
    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text
    assert "Jean Dupont" not in all_text
    assert "PERSON" in all_text
