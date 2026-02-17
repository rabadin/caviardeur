"""Generate test fixture files (.docx, .xlsx, .pptx, .pdf) for integration testing.

Run: uv run python tests/fixtures/create_fixtures.py
"""

from pathlib import Path

import fitz  # PyMuPDF
import openpyxl
from docx import Document
from pptx import Presentation

HERE = Path(__file__).parent


def create_docx():
    doc = Document()
    doc.add_heading("Fiche Client", level=1)
    doc.add_paragraph("Nom: Jean Dupont")
    doc.add_paragraph("Entreprise: Nextech Solutions SAS")
    doc.add_paragraph("Adresse: 12 rue de la Paix, 75002 Paris")
    doc.add_paragraph("SIRET: 732 829 320 00074")
    doc.add_paragraph("")
    doc.add_paragraph("Contact secondaire: Marie Laurent")
    doc.add_paragraph("Adresse: 45 boulevard Haussmann, 75009 Paris")
    doc.save(str(HERE / "sample.docx"))
    print("Created sample.docx")


def create_xlsx():
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Clients"

    headers = ["Nom", "Entreprise", "SIRET", "Adresse", "Contact"]
    for col, h in enumerate(headers, 1):
        ws.cell(row=1, column=col, value=h)

    data = [
        [
            "Jean Dupont",
            "Nextech Solutions SAS",
            "73282932000074",
            "12 rue de la Paix, 75002 Paris",
            "Marie Laurent",
        ],
        [
            "Pierre Martin",
            "DataFlow Industries",
            "44232206200024",
            "8 avenue des Champs-Élysées, 75008 Paris",
            "Sophie Bernard",
        ],
    ]

    for row_idx, row_data in enumerate(data, 2):
        for col_idx, value in enumerate(row_data, 1):
            ws.cell(row=row_idx, column=col_idx, value=value)

    wb.save(str(HERE / "sample.xlsx"))
    print("Created sample.xlsx")


def create_pptx():
    prs = Presentation()

    # Slide 1: title
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Présentation Client"
    slide.placeholders[1].text = "Nextech Solutions SAS"

    # Slide 2: content
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Informations"
    body = slide.placeholders[1]
    body.text = "Client: Jean Dupont"
    body.text_frame.add_paragraph().text = "SIRET: 73282932000074"
    body.text_frame.add_paragraph().text = "Adresse: 12 rue de la Paix, 75002 Paris"
    body.text_frame.add_paragraph().text = "Contact: Marie Laurent"

    prs.save(str(HERE / "sample.pptx"))
    print("Created sample.pptx")


def create_pdf():
    doc = fitz.open()
    page = doc.new_page()
    text = (
        "Fiche Client\n\n"
        "Nom: Jean Dupont\n"
        "Entreprise: Nextech Solutions SAS\n"
        "SIRET: 73282932000074\n"
        "Adresse: 12 rue de la Paix, 75002 Paris\n\n"
        "Contact secondaire: Marie Laurent\n"
    )
    page.insert_text((72, 72), text, fontsize=12, fontname="helv")
    doc.save(str(HERE / "sample.pdf"))
    doc.close()
    print("Created sample.pdf")


if __name__ == "__main__":
    create_docx()
    create_xlsx()
    create_pptx()
    create_pdf()
